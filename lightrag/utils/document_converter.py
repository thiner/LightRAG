"""
Module for document conversion utilities
"""

import pipmaster as pm
from pathlib import Path
from io import BytesIO
from typing import Union, Optional

from lightrag.utils import logger
from lightrag.api.config import args

# set default file size limit to 10MB
if not hasattr(args, 'max_file_size_mb'):
    args.max_file_size_mb = 10

# Set default package names and versions
if not hasattr(args, 'pdf_reader'):
    args.pdf_reader = "pypdf2"
if not hasattr(args, 'pdf_reader_version'):
    args.pdf_reader_version = "latest"

if not hasattr(args, 'docx_reader'):
    args.docx_reader = "python-docx"
if not hasattr(args, 'docx_reader_version'):
    args.docx_reader_version = "latest"
    
if not hasattr(args, 'docx_reader_fallback'):
    args.docx_reader_fallback = "docx"
if not hasattr(args, 'docx_reader_fallback_version'):
    args.docx_reader_fallback_version = "latest"
    
if not hasattr(args, 'pptx_reader'):
    args.pptx_reader = "python-pptx"
if not hasattr(args, 'pptx_reader_version'):
    args.pptx_reader_version = "latest"
    
if not hasattr(args, 'pptx_reader_fallback'):
    args.pptx_reader_fallback = "pptx"
if not hasattr(args, 'pptx_reader_fallback_version'):
    args.pptx_reader_fallback_version = "latest"
    
if not hasattr(args, 'xlsx_reader'):
    args.xlsx_reader = "openpyxl"
if not hasattr(args, 'xlsx_reader_version'):
    args.xlsx_reader_version = "latest"
    
if not hasattr(args, 'doc_converter'):
    args.doc_converter = "docling"
if not hasattr(args, 'doc_converter_version'):
    args.doc_converter_version = "latest"


def install_package_with_version(package_name: str, version: str = "latest"):
    """
    Install a package with the specified version
    
    Args:
        package_name: Name of the package to install
        version: Version specification (default: "latest")
    """
    
    logger.info(f"Installing {package_name} version {version}")
    pm.install_version(package_name, version)


class DocumentConverter:
    """
    Utility class for converting various document formats to text
    """

    def __init__(self, engine: str = "DOCLING"):
        """
        Initialize the DocumentConverter with specified engine
        
        Args:
            engine: Document conversion engine to use ("DOCLING" or "DEFAULT")
        """
        self.engine = engine
        
    def convert_file(self, file_path: Union[str, Path], max_size_mb: int = None) -> str:
        """
        Convert a document file to text content
        
        Args:
            file_path: Path to the file or file name (for extension detection)
            max_size_mb: Maximum file size in MB (default: None, uses global setting)
                
        Returns:
            str: Extracted text content from the document
            
        Raises:
            ValueError: If file is too large or has unsupported extension
        """
        # Check if max_size_mb is provided, otherwise use the global setting    
        if max_size_mb is None:
            max_size_mb = args.max_file_size_mb
            
        if isinstance(file_path, str):
            file_path = Path(file_path)
            
        # Check if file exists
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        # Check file size
        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        if file_size_mb > max_size_mb:
            raise ValueError(f"File size ({file_size_mb:.2f}MB) exceeds maximum allowed size ({max_size_mb}MB)")
        
        # Read file content
        if file_path.suffix.lower() in ('.txt', '.md', '.html', '.htm', '.py', '.js'):
            # For text files, try to open directly
            file_content = None  # Will be read in _convert_text_file
        else:
            # For binary files, read content first
            with open(file_path, "rb") as f:
                file_content = f.read()
                
        ext = file_path.suffix.lower()
        
        # Process based on file type
        if ext in ('.txt', '.md', '.html', '.htm', '.tex', '.json', '.xml', 
                  '.yaml', '.yml', '.rtf', '.odt', '.epub', '.csv', '.log', 
                  '.conf', '.ini', '.properties', '.sql', '.bat', '.sh', 
                  '.c', '.cpp', '.py', '.java', '.js', '.ts', '.swift', 
                  '.go', '.rb', '.php', '.css', '.scss', '.less'):
            return self._convert_text_file(file_path, file_content)
        elif ext == '.pdf':
            return self._convert_pdf(file_path, file_content)
        elif ext == '.docx':
            return self._convert_docx(file_path, file_content)
        elif ext == '.pptx':
            return self._convert_pptx(file_path, file_content)
        elif ext == '.xlsx':
            return self._convert_xlsx(file_path, file_content)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    def _convert_text_file(self, file_path: Path, file_content: Optional[bytes]) -> str:
        """Convert plain text files"""
        if file_content:
            try:
                return file_content.decode("utf-8")
            except UnicodeDecodeError:
                logger.error(f"File {file_path.name} is not valid UTF-8 encoded text.")
                raise
        else:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
    
    def _convert_pdf(self, file_path: Path, file_content: Optional[bytes]) -> str:
        """Convert PDF files"""
        if self.engine == "DOCLING":
            return self._convert_with_docling(file_path)
        else:
            install_package_with_version(args.pdf_reader, args.pdf_reader_version)
            from pypdf2 import PdfReader
            
            content = ""
            if file_content:
                pdf_file = BytesIO(file_content)
                reader = PdfReader(pdf_file)
            else:
                reader = PdfReader(str(file_path))
                
            for page in reader.pages:
                content += page.extract_text() + "\n"
            return content
    
    def _convert_docx(self, file_path: Path, file_content: Optional[bytes]) -> str:
        """Convert DOCX files"""
        if self.engine == "DOCLING":
            return self._convert_with_docling(file_path)
        else:
            try:
                install_package_with_version(args.docx_reader, args.docx_reader_version)
            except Exception:
                install_package_with_version(args.docx_reader_fallback, args.docx_reader_fallback_version)
            from docx import Document
            
            if file_content:
                docx_file = BytesIO(file_content)
                doc = Document(docx_file)
            else:
                doc = Document(file_path)
                
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _convert_pptx(self, file_path: Path, file_content: Optional[bytes]) -> str:
        """Convert PPTX files"""
        if self.engine == "DOCLING":
            return self._convert_with_docling(file_path)
        else:
            try:
                install_package_with_version(args.pptx_reader, args.pptx_reader_version)
            except Exception:
                install_package_with_version(args.pptx_reader_fallback, args.pptx_reader_fallback_version)
            from pptx import Presentation
            
            content = ""
            if file_content:
                pptx_file = BytesIO(file_content)
                prs = Presentation(pptx_file)
            else:
                prs = Presentation(file_path)
                
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content
    
    def _convert_xlsx(self, file_path: Path, file_content: Optional[bytes]) -> str:
        """Convert XLSX files"""
        if self.engine == "DOCLING":
            return self._convert_with_docling(file_path)
        else:
            install_package_with_version(args.xlsx_reader, args.xlsx_reader_version)
            from openpyxl import load_workbook
            
            content = ""
            if file_content:
                xlsx_file = BytesIO(file_content)
                wb = load_workbook(xlsx_file)
            else:
                wb = load_workbook(file_path)
                
            for sheet in wb:
                content += f"Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    content += "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    ) + "\n"
                content += "\n"
            return content
    
    def _convert_with_docling(self, file_path: Path) -> str:
        """Convert file using the docling library"""
        install_package_with_version(args.doc_converter, args.doc_converter_version)
        from docling.document_converter import DocumentConverter
        
        converter = DocumentConverter()
        result = converter.convert(file_path)
        return result.document.export_to_markdown()